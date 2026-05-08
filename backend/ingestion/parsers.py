from __future__ import annotations

from pathlib import Path

def parse_pdf(path: Path) -> list[str]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    return [(page.extract_text() or "").strip() for page in reader.pages]


def parse_docx(path: Path) -> list[str]:
    from docx import Document

    document = Document(str(path))
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs]
    return [paragraph for paragraph in paragraphs if paragraph]


def _shape_texts(slide: object) -> list[str]:
    texts: list[str] = []
    for shape in getattr(slide, "shapes", []):
        if hasattr(shape, "text_frame") and shape.text_frame:
            text = shape.text_frame.text.strip()
            if text:
                texts.append(text)
    return texts


def _notes_text(slide: object) -> str:
    try:
        notes_slide = slide.notes_slide
    except Exception:
        return ""

    notes_frame = getattr(notes_slide, "notes_text_frame", None)
    if notes_frame is None:
        return ""
    return notes_frame.text.strip()


def parse_pptx(path: Path) -> list[str]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    slides: list[str] = []
    for slide in presentation.slides:
        parts = _shape_texts(slide)
        notes = _notes_text(slide)
        if notes:
            parts.append(notes)
        slides.append("\n\n".join(parts).strip())
    return slides


def parse_text(path: Path) -> list[str]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    return [part.strip() for part in text.split("\n\n") if part.strip()]
